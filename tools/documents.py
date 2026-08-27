"""
tools/documents.py

Renders NAVI's plain markdown-ish output into PDF, DOCX, or PPTX bytes —
used when a command's opt-in file request asks for a specific format
instead of the default .md text save.

All three libraries (fpdf2, python-docx, python-pptx) are pure Python
with no system-level binary dependency (unlike e.g. weasyprint, which
needs Cairo/Pango installed) — a deliberate fit for Render's constrained
environment, matching the "no heavy SDKs" rule the rest of NAVI's
dependencies already follow.

Unicode note: fpdf2's built-in core fonts (Helvetica, Times, Courier)
are Latin-1 only and will mangle or crash on the em-dashes, accented
characters, and smart quotes this content routinely contains — real
risk, not theoretical (see the ollama/ollama#17836 investigation
earlier in this project's history, a different but related class of
multi-byte-character bug). Rather than bundle a new font file in the
repo, this reuses the DejaVu Sans TTF matplotlib already ships as a
dependency (NAVI already requires matplotlib for chart rendering) —
genuinely Unicode-capable, zero new binary assets to check in or
license.
"""

import re
from io import BytesIO

import matplotlib
from docx import Document
from docx.shared import Pt as DocxPt
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt

_DEJAVU_PATH = f"{matplotlib.get_data_path()}/fonts/ttf/DejaVuSans.ttf"
_DEJAVU_BOLD_PATH = f"{matplotlib.get_data_path()}/fonts/ttf/DejaVuSans-Bold.ttf"


class DocumentRenderError(Exception):
    pass


def _heading_level(line: str) -> tuple[int, str]:
    """Returns (level, text) for a markdown heading line, or (0, line)
    if it isn't one. Only # through ### matter here — NAVI's own output
    doesn't go deeper than that in practice."""
    match = re.match(r"^(#{1,3})\s+(.*)", line)
    if not match:
        return 0, line
    return len(match.group(1)), match.group(2)


def render_pdf(title: str, content: str) -> bytes:
    """Minimal markdown-aware PDF: '#'/'##'/'###' lines become headings
    at decreasing size, everything else is body text. Not full markdown
    rendering (no bold/italic/links/tables) — just enough structure to
    read like a real document instead of a wall of undifferentiated
    text, matching what NAVI's own commands actually produce."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_font("DejaVu", "", _DEJAVU_PATH)
    pdf.add_font("DejaVu", "B", _DEJAVU_BOLD_PATH)

    pdf.set_font("DejaVu", "B", 18)
    pdf.multi_cell(0, 10, title)
    pdf.ln(4)

    heading_sizes = {1: 15, 2: 13, 3: 12}
    for raw_line in content.split("\n"):
        line = raw_line.rstrip()
        if not line.strip():
            pdf.ln(3)
            continue
        level, text = _heading_level(line)
        if level:
            pdf.set_font("DejaVu", "B", heading_sizes.get(level, 12))
            pdf.multi_cell(0, 8, text)
        else:
            pdf.set_font("DejaVu", "", 11)
            pdf.multi_cell(0, 6, line)
        # multi_cell leaves the cursor x near the right margin, not
        # reset to the left — without this, two consecutive non-blank
        # lines (e.g. a heading directly followed by body text with no
        # blank line between them) crash the next multi_cell call with
        # "not enough horizontal space," since it's measuring from
        # wherever x already is, not from the left margin.
        pdf.ln()

    return bytes(pdf.output())


def render_docx(title: str, content: str) -> bytes:
    """Same heading-aware structure as render_pdf, using Word's real
    heading styles rather than just bigger/bold text — makes the
    document actually navigable (table of contents, outline view) in
    Word, not just visually styled."""
    doc = Document()
    doc.add_heading(title, level=0)

    for raw_line in content.split("\n"):
        line = raw_line.rstrip()
        if not line.strip():
            continue
        level, text = _heading_level(line)
        if level:
            doc.add_heading(text, level=level)
        else:
            doc.add_paragraph(line)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def render_pptx(title: str, content: str) -> bytes:
    """Turns '##'-level headings into new slides, with everything under
    each one as bullet points on that slide — a reasonable mapping from
    a structured document into a deck, not a literal 1:1 conversion
    (PPTX isn't really a natural fit for NAVI's prose-heavy output, but
    this gives a usable starting deck rather than one slide of a wall
    of text)."""
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title

    current_slide = None
    current_body = None
    for raw_line in content.split("\n"):
        line = raw_line.rstrip()
        if not line.strip():
            continue
        level, text = _heading_level(line)
        if level and level <= 2:
            current_slide = prs.slides.add_slide(bullet_layout)
            current_slide.shapes.title.text = text
            current_body = current_slide.placeholders[1].text_frame
            current_body.clear()
        elif current_body is not None:
            p = current_body.paragraphs[0] if not current_body.paragraphs[0].text else current_body.add_paragraph()
            p.text = text if level else line
        # Content before any heading has nowhere to go on a bullet
        # slide — dropped rather than forced onto the title slide.

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


RENDERERS = {
    "pdf": render_pdf,
    "docx": render_docx,
    "pptx": render_pptx,
}


def render(format_name: str, title: str, content: str) -> bytes:
    renderer = RENDERERS.get(format_name.lower())
    if not renderer:
        raise DocumentRenderError(f"Unsupported format: {format_name}")
    return renderer(title, content)
